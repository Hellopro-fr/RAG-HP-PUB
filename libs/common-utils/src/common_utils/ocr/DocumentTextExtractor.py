import os
import sys
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Union
import logging
import urllib.parse
from urllib.error import URLError, HTTPError
import mimetypes
from tempfile import NamedTemporaryFile

from common_utils.ocr.OCRDocExtractor import OCRDocExtractor

# Import des bibliothèques nécessaires
try:
    from PIL import Image
    import fitz  # PyMuPDF
    from docx import Document
    import openpyxl
    from pptx import Presentation
    import mammoth
    from odf.opendocument import load
    from odf.text import P
    from odf.element import Text
    import subprocess
    import tempfile
    import requests
except ImportError as e:
    print(f"Erreur d'import: {e}")
    print("Installez les dépendances manquantes avec:")
    print("pip install Pillow PyMuPDF python-docx openpyxl python-pptx mammoth odfpy requests")
    sys.exit(1)

class DocumentTextExtractor:
    """
    Classe pour extraire le texte de différents types de documents
    """
    
    def __init__(self, download_dir: Optional[str] = None, auto_cleanup: bool = True):
        """
        Initialise l'extracteur
        
        Args:
            download_dir: Répertoire pour télécharger les fichiers depuis les URLs
            auto_cleanup: Supprimer automatiquement les fichiers après traitement
        """
        # Répertoire de téléchargement
        self.download_dir = Path(download_dir) if download_dir else Path("./downloaded_files")
        self.download_dir.mkdir(exist_ok=True)
        
        # Configuration du nettoyage automatique
        self.auto_cleanup = auto_cleanup
        self.files_to_cleanup = []  # Liste des fichiers à supprimer
        
        # Initialiser OCRExtractor (une seule fois)
        self.ocr_processor = OCRDocExtractor()
        
        # Formats supportés
        self.image_formats = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
        self.ocr_supported = {'.png', '.bmp', '.pdf', '.gif', '.jpg', '.jpeg'}
        self.document_formats = {'.doc', '.docx', '.xlsx', '.xls', '.pptx', '.ppt', '.odt'}
        
        # Configuration du logging
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)
    
    def is_url(self, path_or_url: str) -> bool:
        """
        Vérifie si une chaîne est une URL
        
        Args:
            path_or_url: Chaîne à vérifier
            
        Returns:
            True si c'est une URL
        """
        return path_or_url.startswith(('http://', 'https://'))
    
    def get_filename_from_url(self, url: str) -> str:
        """
        Extrait le nom de fichier depuis une URL
        
        Args:
            url: URL du fichier
            
        Returns:
            Nom de fichier
        """
        # Parser l'URL
        parsed_url = urllib.parse.urlparse(url)
        filename = os.path.basename(parsed_url.path)
        
        # Si pas de nom de fichier dans l'URL, générer un nom
        if not filename or '.' not in filename:
            # Essayer de deviner l'extension depuis le Content-Type
            try:
                response = requests.head(url, timeout=10)
                content_type = response.headers.get('Content-Type', '')
                ext = mimetypes.guess_extension(content_type.split(';')[0])
                filename = f"downloaded_file{ext if ext else '.bin'}"
            except:
                filename = "downloaded_file.bin"
        
        return filename
    
    def download_file(self, url: str) -> Path:
        """
        Télécharge un fichier depuis une URL
        
        Args:
            url: URL du fichier à télécharger
            
        Returns:
            Chemin vers le fichier téléchargé
        """
        try:
            self.logger.info(f"Téléchargement de: {url}")
            
            filename = self.get_filename_from_url(url)
            local_path = self.download_dir / filename
            
            # Éviter les téléchargements en double
            if local_path.exists():
                self.logger.info(f"Fichier déjà téléchargé: {local_path}")
                self.add_file_for_cleanup(local_path, is_downloaded=True)
                return local_path
            
            # Configuration des headers pour éviter les blocages
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(url, headers=headers, timeout=30, stream=True)
            response.raise_for_status()
            
            # Vérifier la taille du fichier (limite de sécurité: 100MB)
            content_length = response.headers.get('Content-Length')
            if content_length and int(content_length) > 100 * 1024 * 1024:
                raise Exception(f"Fichier trop volumineux: {content_length} bytes")
            
            # Télécharger le fichier
            with open(local_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            
            self.logger.info(f"Fichier téléchargé: {local_path}")
            self.add_file_for_cleanup(local_path, is_downloaded=True)
            return local_path
            
        except Exception as e:
            self.logger.error(f"Erreur lors du téléchargement de {url}: {e}")
            raise
    
    def resolve_path_or_url(self, path_or_url: Union[str, Path]) -> Path:
        """
        Résout un chemin ou URL vers un chemin local
        
        Args:
            path_or_url: Chemin local ou URL
            
        Returns:
            Chemin local vers le fichier
        """
        if isinstance(path_or_url, str) and self.is_url(path_or_url):
            return self.download_file(path_or_url)
        else:
            return Path(path_or_url)
        
    def convert_image_to_supported_format(self, image_path: Path) -> Path:
        """
        Convertit une image vers un format supporté par l'OCR
        
        Args:
            image_path: Chemin vers l'image source
            
        Returns:
            Chemin vers l'image convertie
        """
        try:
            with Image.open(image_path) as img:
                # Convertir en RGB si nécessaire
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                
                # Nouveau nom de fichier avec extension .png
                output_path = image_path.parent / f"{image_path.stem}_converted.png"
                img.save(output_path, 'PNG')
                
                self.logger.info(f"Image convertie: {image_path} -> {output_path}")
                return output_path
                
        except Exception as e:
            self.logger.error(f"Erreur lors de la conversion de {image_path}: {e}")
            raise
    
    def add_file_for_cleanup(self, file_path: Path, is_downloaded: bool = False):
        """
        Ajoute un fichier à la liste de nettoyage
        
        Args:
            file_path: Chemin du fichier
            is_downloaded: Si le fichier a été téléchargé
        """
        if self.auto_cleanup:
            self.files_to_cleanup.append({
                'path': file_path,
                'is_downloaded': is_downloaded
            })
    
    def cleanup_files(self):
        """
        Supprime tous les fichiers marqués pour le nettoyage
        """
        for file_info in self.files_to_cleanup:
            try:
                file_path = file_info['path']
                if file_path.exists():
                    file_path.unlink()
                    self.logger.info(f"Fichier supprimé: {file_path}")
            except Exception as e:
                self.logger.warning(f"Impossible de supprimer {file_path}: {e}")
        
        self.files_to_cleanup.clear()

    def extract_text_from_image_ocr(self, file_path: Path) -> str:
        """
        Extrait le texte d'une image ou d'un PDF avec OCRExtractor.
        - Si un PDF contient plus de 10 pages, il est traité par blocs de 4 pages.
        """
        # --- Vérifications de base ---
        if not file_path.exists():
            raise FileNotFoundError(f"Fichier introuvable : {file_path}")

        if not file_path.is_file():
            raise ValueError(f"Le chemin fourni n'est pas un fichier valide : {file_path}")

        ext = file_path.suffix.lower()

        if ext not in self.ocr_supported:
            raise ValueError(f"Format de fichier non supporté : {ext}")

        # --- Si PDF ---
        if ext == ".pdf":
            try:
                doc = fitz.open(file_path)
                num_pages = doc.page_count
            except Exception as e:
                raise RuntimeError(f"Impossible d'ouvrir le PDF : {e}")

            # --- Si plus de 10 pages, traitement par blocs de 4 pages ---
            if num_pages > 10:
                all_text = []
                block_size = 4

                for start in range(0, num_pages, block_size):
                    end = min(start + block_size - 1, num_pages - 1)

                    # Nouveau document temporaire contenant 4 pages
                    block_doc = fitz.open()
                    block_doc.insert_pdf(doc, from_page=start, to_page=end)

                    # Créer le fichier temporaire
                    with NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
                        block_doc.save(temp_pdf.name)
                        temp_path = Path(temp_pdf.name)

                    # OCR sur le bloc de 4 pages
                    try:
                        text = self.ocr_processor.convert_doc_to_markdown([temp_path])
                        all_text.append(text)
                    finally:
                        temp_path.unlink(missing_ok=True)  # suppression du fichier temporaire

                doc.close()
                return "\n".join(all_text)

            doc.close()

        # --- Sinon : image ou petit PDF ---
        return self.ocr_processor.convert_doc_to_markdown([file_path])
        
    def has_extractable_images(self, file_path: Path) -> bool:
        """
        Vérifie si un document contient des images non extractibles
        
        Args:
            file_path: Chemin vers le document
            
        Returns:
            True si le document contient des images difficiles à extraire
        """
        try:
            if file_path.suffix.lower() in {'.ppt', '.pptx'}:
                # PowerPoint contient souvent des images intégrées
                return True
            elif file_path.suffix.lower() in {'.doc','.docx'}:
                # Les anciens formats Word peuvent avoir des images intégrées
                return True
            elif file_path.suffix.lower() in {'.xls','.xlsx'}:
                # Excel ancien format peut avoir des images
                return True
            
            return False
            
        except Exception:
            return True  # En cas de doute, convertir en PDF
    
    def convert_to_pdf(self, file_path: Path) -> Path:
        """
        Convertit un document en PDF en utilisant LibreOffice
        
        Args:
            file_path: Chemin vers le document source
            
        Returns:
            Chemin vers le PDF généré
        """
        try:
            output_dir = file_path.parent
            pdf_path = output_dir / f"{file_path.stem}.pdf"
            
            # Commande LibreOffice pour conversion
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(output_dir),
                str(file_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0 and pdf_path.exists():
                self.logger.info(f"Document converti en PDF: {file_path} -> {pdf_path}")
                return pdf_path
            else:
                raise Exception(f"Échec de la conversion: {result.stderr}")
                
        except Exception as e:
            self.logger.error(f"Erreur lors de la conversion PDF de {file_path}: {e}")
            raise
    
    def extract_text_from_pdf(self, pdf_path: Path, use_ocr: bool = False) -> str:
        """
        Extrait le texte d'un PDF
        
        Args:
            pdf_path: Chemin vers le PDF
            use_ocr: Utiliser l'OCR si le texte n'est pas extractible
            
        Returns:
            Texte extrait
        """
        try:
            text = ""

            if use_ocr:
                text = self.extract_text_from_image_ocr(pdf_path)
            else:
                doc = fitz.open(pdf_path)
                
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    page_text = page.get_text()
                    
                    if not page_text.strip() and use_ocr:
                        # Si pas de texte extractible, utiliser OCR
                        pix = page.get_pixmap()
                        img_data = pix.tobytes("png")
                        
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                            tmp.write(img_data)
                            tmp_path = Path(tmp.name)
                        
                        try:
                            page_text = self.extract_text_from_image_ocr(tmp_path)
                        finally:
                            tmp_path.unlink(missing_ok=True)
                    
                    text += page_text + "\n"
                
                doc.close()
            
            return text.strip()
            
        except Exception as e:
            self.logger.error(f"Erreur lors de l'extraction PDF {pdf_path}: {e}")
            return ""
    
    def extract_text_from_docx(self, docx_path: Path) -> str:
        """Extrait le texte d'un fichier DOCX"""
        try:
            doc = Document(docx_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text.strip()
        except Exception as e:
            self.logger.error(f"Erreur DOCX {docx_path}: {e}")
            return ""
    
    def extract_text_from_doc(self, doc_path: Path) -> str:
        """Extrait le texte d'un fichier DOC avec mammoth"""
        try:
            with open(doc_path, "rb") as docfile:
                result = mammoth.extract_raw_text(docfile)
                return result.value.strip()
        except Exception as e:
            self.logger.error(f"Erreur DOC {doc_path}: {e}")
            return ""
    
    def extract_text_from_xlsx(self, xlsx_path: Path) -> str:
        """Extrait le texte d'un fichier XLSX"""
        try:
            wb = openpyxl.load_workbook(xlsx_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"=== Feuille: {sheet_name} ===")
                
                for row in ws.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        text_parts.append(" | ".join(row_text))
            
            wb.close()
            return "\n".join(text_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur XLSX {xlsx_path}: {e}")
            return ""
    
    def extract_text_from_pptx(self, pptx_path: Path) -> str:
        """Extrait le texte d'un fichier PPTX"""
        try:
            prs = Presentation(pptx_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== Slide {slide_num} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur PPTX {pptx_path}: {e}")
            return ""
    
    def extract_text_from_odt(self, odt_path: Path) -> str:
        """Extrait le texte d'un fichier ODT"""
        try:
            doc = load(odt_path)
            text_parts = []
            
            for paragraph in doc.getElementsByType(P):
                text_content = []
                for node in paragraph.childNodes:
                    if node.nodeType == node.TEXT_NODE:
                        text_content.append(str(node.data))
                    elif hasattr(node, 'data'):
                        text_content.append(str(node.data))
                
                if text_content:
                    text_parts.append(''.join(text_content))
            
            return "\n".join(text_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur ODT {odt_path}: {e}")
            return ""
    
    def process_single_file(self, path_or_url: Union[str, Path]) -> Dict[str, str]:
        """
        Traite un seul fichier (local ou URL) et extrait son texte
        
        Args:
            path_or_url: Chemin vers le fichier local ou URL
            
        Returns:
            Dictionnaire avec les informations d'extraction
        """
        result = {
            'original_path': str(path_or_url),
            'file_path': '',
            'file_type': '',
            'text': '',
            'method': '',
            'status': 'success',
            'error': None,
            'is_url': False
        }
        
        try:
            # Résoudre le chemin ou URL
            if isinstance(path_or_url, str) and self.is_url(path_or_url):
                result['is_url'] = True
                file_path = self.resolve_path_or_url(path_or_url)
                result['method'] = f"Téléchargé depuis URL -> "
            else:
                file_path = Path(path_or_url)
            
            result['file_path'] = str(file_path)
            result['file_type'] = file_path.suffix.lower()
            
            file_ext = file_path.suffix.lower()
            
            # Traitement des images
            if file_ext in self.image_formats:
                if file_ext not in self.ocr_supported:
                    converted_path = self.convert_image_to_supported_format(file_path)
                    self.add_file_for_cleanup(converted_path)
                    result['text'] = self.extract_text_from_image_ocr(converted_path)
                else:
                    result['text'] = self.extract_text_from_image_ocr(file_path)
                result['method'] += 'OCRExtractor'
            
            # Traitement des PDFs
            elif file_ext == '.pdf':
                result['text'] = self.extract_text_from_pdf(file_path,use_ocr=True)
                result['method'] += 'PDF extraction + OCR'
            
            # Traitement des documents avec vérification d'images
            elif file_ext in self.document_formats:
                if self.has_extractable_images(file_path):
                    # Convertir en PDF puis utiliser OCR
                    pdf_path = self.convert_to_pdf(file_path)
                    self.add_file_for_cleanup(pdf_path)
                    result['text'] = self.extract_text_from_pdf(pdf_path)
                    result['method'] += 'Conversion PDF + PDF extraction'
                    # Optionnel: supprimer le PDF temporaire
                    pdf_path.unlink(missing_ok=True)
                else:
                    # Extraction directe
                    if file_ext == '.docx':
                        result['text'] = self.extract_text_from_docx(file_path)
                        result['method'] += 'DOCX direct'
                    elif file_ext == '.doc':
                        result['text'] = self.extract_text_from_doc(file_path)
                        result['method'] += 'DOC direct'
                    elif file_ext in {'.xlsx', '.xls'}:
                        if file_ext == '.xlsx':
                            result['text'] = self.extract_text_from_xlsx(file_path)
                            result['method'] += 'XLSX direct'
                        else:
                            # XLS -> conversion PDF
                            pdf_path = self.convert_to_pdf(file_path)
                            self.add_file_for_cleanup(pdf_path)
                            result['text'] = self.extract_text_from_pdf(pdf_path)
                            result['method'] += 'XLS -> PDF'
                    elif file_ext in {'.pptx', '.ppt'}:
                        if file_ext == '.pptx':
                            result['text'] = self.extract_text_from_pptx(file_path)
                            result['method'] += 'PPTX direct'
                        else:
                            # PPT -> conversion PDF
                            pdf_path = self.convert_to_pdf(file_path)
                            self.add_file_for_cleanup(pdf_path)
                            result['text'] = self.extract_text_from_pdf(pdf_path)
                            result['method'] += 'PPT -> PDF'
                    elif file_ext == '.odt':
                        result['text'] = self.extract_text_from_odt(file_path)
                        result['method'] += 'ODT direct'
            
            else:
                result['status'] = 'unsupported'
                result['error'] = f"Format non supporté: {file_ext}"
        
        except Exception as e:
            result['status'] = 'error'
            result['error'] = str(e)
            self.logger.error(f"Erreur lors du traitement de {path_or_url}: {e}")
        finally:
            # Nettoyage automatique après traitement
            self.cleanup_files()
        
        return result
    
    def process_multiple_files(self, paths_or_urls: List[Union[str, Path]]) -> List[Dict[str, str]]:
        """
        Traite une liste de fichiers/URLs
        
        Args:
            paths_or_urls: Liste de chemins locaux et/ou URLs
            
        Returns:
            Liste des résultats d'extraction
        """
        results = []
        
        for i, path_or_url in enumerate(paths_or_urls, 1):
            print(f"Traitement {i}/{len(paths_or_urls)}: {path_or_url}")
            
            result = self.process_single_file(path_or_url)
            results.append(result)
            
            # Affichage du progrès
            status_symbol = "✓" if result['status'] == 'success' else "✗"
            file_name = Path(result['file_path']).name if result['file_path'] else "Unknown"
            print(f"{status_symbol} {file_name} - {result['method']}")
            
            if result['status'] == 'error':
                print(f"  Erreur: {result['error']}")
        
        return results